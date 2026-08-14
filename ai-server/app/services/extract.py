"""Content extraction for PDF / PPTX / DOCX sources."""
import io
import logging

import httpx

from app.config import get_settings

logger = logging.getLogger(__name__)


async def download(url: str) -> bytes:
    settings = get_settings()
    limit = settings.max_download_mb * 1024 * 1024
    async with httpx.AsyncClient(timeout=45.0, follow_redirects=True) as client:
        resp = await client.get(url)
        resp.raise_for_status()
        if len(resp.content) > limit:
            raise ValueError(f"File exceeds {settings.max_download_mb}MB limit")
        return resp.content


def detect_kind(name: str | None, mime: str | None) -> str:
    blob = f"{name or ''} {mime or ''}".lower()
    if "pdf" in blob:
        return "pdf"
    if "presentation" in blob or blob.endswith(".ppt") or ".pptx" in blob:
        return "ppt"
    if "word" in blob or ".doc" in blob:
        return "doc"
    return "unknown"


def extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages[:80])


def extract_ppt(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    chunks.append(text)
    return "\n".join(chunks)


def extract_doc(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


async def extract_file(url: str, name: str | None, mime: str | None, kind: str) -> tuple[str, str | None]:
    """Returns (text, warning)."""
    resolved = kind if kind != "auto" else detect_kind(name, mime)
    try:
        data = await download(url)
        if resolved == "pdf":
            return extract_pdf(data), None
        if resolved == "ppt":
            return extract_ppt(data), None
        if resolved == "doc":
            return extract_doc(data), None
        return "", f"Unsupported file type for {name or url}"
    except Exception as exc:  # noqa: BLE001 - never fail the whole request on one file
        logger.warning("extraction failed for %s: %s", url, exc)
        return "", f"Could not read {name or url}: {exc}"